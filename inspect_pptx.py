from pptx import Presentation
prs = Presentation(r"C:\Users\gagal\Downloads\template vídeo - trabalho MLP - 2026.pptx")
for i, slide in enumerate(prs.slides):
    print(f"=== Slide {i+1} ({len(slide.shapes)} shapes) ===")
    for j, sh in enumerate(slide.shapes):
        pos = f"L={sh.left//9144:4}pt T={sh.top//9144:4}pt W={sh.width//9144:4}pt H={sh.height//9144:4}pt"
        txt = sh.text_frame.text[:70].replace("\n"," ") if sh.has_text_frame else "(no text)"
        print(f"  [{j}] type={sh.shape_type} {pos}  \"{txt}\"")
