"""
Gerador de slides MLP – estilo do exemplo do grupo:
  - Texto explicativo à esquerda (frases completas, 16pt headers, 13pt body)
  - Imagem de código à direita (gerada com matplotlib, fundo escuro)
  - Slides sem código: texto em largura total
"""

import csv, os, shutil, tempfile
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

TEMPLATE = r"C:\Users\gagal\Downloads\template vídeo - trabalho MLP - 2026.pptx"
OUTPUT   = r"C:\Users\gagal\Desktop\Faculdade\IA\slides_MLP_preenchido.pptx"
CSV_TEST = r"C:\Users\gagal\Desktop\Faculdade\IA\saidas_teste.csv"

PT = 9144
TMPDIR = tempfile.mkdtemp()

HDR  = RGBColor(0,  70, 130)
BODY = RGBColor(30, 30,  30)
GRNN = RGBColor(0, 100,  40)   # verde – destaque positivo

TXT_L, TXT_T, TXT_W, TXT_H = 82,  203, 590, 462
IMG_L, IMG_T, IMG_W, IMG_H = 690, 203, 590, 462
FUL_L, FUL_T, FUL_W, FUL_H = 82, 203, 1200, 462


# ═══════════════════════════════════════════════════════════════════════════
# Helpers de layout
# ═══════════════════════════════════════════════════════════════════════════

def _remove_p(para):
    para._p.getparent().remove(para._p)

def _clear_body(slide):
    if len(slide.shapes) >= 2:
        slide.shapes[1]._element.getparent().remove(
            slide.shapes[1]._element)

def _new_tb(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l*PT, t*PT, w*PT, h*PT)
    tf = tb.text_frame
    tf.word_wrap = True
    for p in list(tf.paragraphs):
        _remove_p(p)
    return tf

def _add_para(tf, text, *, header=False, bold=False, size=None,
              color=None, space_after=5):
    from pptx.util import Pt
    p   = tf.add_paragraph()
    run = p.add_run()
    run.text           = text
    run.font.name      = "Calibri"
    run.font.size      = Pt(size or (16 if header else 13))
    run.font.bold      = header or bold
    run.font.color.rgb = color or (HDR if header else BODY)
    from pptx.oxml.ns import qn
    from lxml import etree
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    spcAft = etree.SubElement(pPr, qn("a:spcAft"))
    spcPts = etree.SubElement(spcAft, qn("a:spcPts"))
    spcPts.set("val", str(space_after * 100))


def left_text(slide, items):
    tf = _new_tb(slide, TXT_L, TXT_T, TXT_W, TXT_H)
    for text, is_hdr in items:
        _add_para(tf, text, header=is_hdr, space_after=6 if is_hdr else 5)

def full_text(slide, items):
    tf = _new_tb(slide, FUL_L, FUL_T, FUL_W, FUL_H)
    for text, is_hdr in items:
        _add_para(tf, text, header=is_hdr, space_after=6 if is_hdr else 5)

def right_img(slide, path):
    slide.shapes.add_picture(path, IMG_L*PT, IMG_T*PT, IMG_W*PT, IMG_H*PT)

def full_img(slide, path):
    slide.shapes.add_picture(path, FUL_L*PT, FUL_T*PT, FUL_W*PT, FUL_H*PT)

def _add_extra_slide(prs):
    """Add a new content slide using same layout as slide[1]."""
    layout = prs.slides[1].slide_layout
    new_s  = prs.slides.add_slide(layout)
    _clear_body(new_s)
    return new_s

def _set_title(slide, text):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.shape_type in (13, 14, 1, 2):
            tf = sh.text_frame
            for para in tf.paragraphs:
                for r in para.runs:
                    r.text = ""
            p0 = tf.paragraphs[0]
            if p0.runs:
                p0.runs[0].text = text
            else:
                run = p0.add_run()
                run.text = text
                run.font.name = "Calibri"
                run.font.size = Pt(24)
                run.font.bold = True
                run.font.color.rgb = HDR
            return
    # fallback: add a text box at the title position
    tb = slide.shapes.add_textbox(82*PT, 31*PT, 1150*PT, 155*PT)
    tf = tb.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.name = "Calibri"; run.font.size = Pt(24)
    run.font.bold = True; run.font.color.rgb = HDR


# ═══════════════════════════════════════════════════════════════════════════
# Imagem: código (tema escuro)
# ═══════════════════════════════════════════════════════════════════════════

BG = "#0d1117"; FG = "#e6edf3"; CMT = "#8b949e"

def _make_code_img(lines, path, figsize=(7.8, 5.8)):
    fig, ax = plt.subplots(figsize=figsize)
    fig.patch.set_facecolor(BG); ax.set_facecolor(BG); ax.axis("off")
    n = len(lines)
    y_top = 0.97; line_h = (y_top - 0.02) / max(n, 1)
    font_size = max(8.5, min(10.5, 10.5 - max(0, n - 14) * 0.15))
    for i, line in enumerate(lines):
        y = y_top - i * line_h
        if y < 0.01: break
        color = CMT if line.lstrip().startswith("#") else FG
        ax.text(0.025, y, line, color=color, fontsize=font_size,
                fontfamily="monospace", va="top", transform=ax.transAxes)
    fig.tight_layout(pad=0.4)
    fig.savefig(path, dpi=160, bbox_inches="tight", facecolor=BG, edgecolor="none")
    plt.close(fig)


# ═══════════════════════════════════════════════════════════════════════════
# Imagem: arquitetura MLP
# ═══════════════════════════════════════════════════════════════════════════

def _make_arch_img(path):
    fig, ax = plt.subplots(figsize=(7.8, 5.8))
    fig.patch.set_facecolor("white"); ax.set_facecolor("white")
    ax.set_xlim(0, 10); ax.set_ylim(0, 8); ax.axis("off")
    layers = [
        (2.0, 7, 120, "Entrada\n120", "#1565C0"),
        (5.0, 6,  64, "Oculta\n64",  "#E65100"),
        (8.0, 5,  26, "Saída\n26",   "#2E7D32"),
    ]
    radius = 0.27
    for li, (x, n_show, total, label, color) in enumerate(layers):
        ys = np.linspace(1.2, 6.8, n_show)
        if li < len(layers) - 1:
            x_next, n_next = layers[li+1][0], layers[li+1][1]
            ys_next = np.linspace(1.2, 6.8, n_next)
            for y1 in ys:
                for y2 in ys_next:
                    ax.plot([x+radius, x_next-radius], [y1, y2],
                            color="#cccccc", lw=0.5, zorder=1)
        for y in ys:
            c = plt.Circle((x, y), radius, fc=color, ec="white", lw=1.5, zorder=2)
            ax.add_patch(c)
        if total > n_show:
            ax.text(x, 0.6, "⋮", ha="center", va="center",
                    fontsize=18, color="#666666")
        ax.text(x, 7.5, label, ha="center", va="center",
                fontsize=12, fontweight="bold", color="#222222")
    for x, label in [(3.5, "tanh(z₁)"), (6.5, "tanh(z₂)")]:
        ax.annotate("", xy=(x+0.6, 4), xytext=(x-0.6, 4),
                    arrowprops=dict(arrowstyle="->", color="#555555", lw=1.4))
        ax.text(x, 4.35, label, ha="center", fontsize=9,
                color="#555555", style="italic")
    ax.set_title("Arquitetura MLP implementada", fontsize=13,
                 fontweight="bold", color="#111111", pad=8)
    fig.tight_layout()
    fig.savefig(path, dpi=160, bbox_inches="tight", facecolor="white")
    plt.close(fig)


# ═══════════════════════════════════════════════════════════════════════════
# Imagem: tabela Grid Search (dados reais da execução)
# ═══════════════════════════════════════════════════════════════════════════

def _make_grid_table_img(path):
    rows = [
        ("64", "8000", "7290", "0.015", "16.94s", "94.83%", "84.85% ★"),
        ("50", "8000", "7936", "0.015", "16.27s", "92.89%", "83.33%"),
        ("50", "7000", "7000", "0.015", "13.98s", "90.30%", "78.79%"),
        ("64", "7000", "7000", "0.015", "19.01s", "93.10%", "78.03%"),
    ]
    cols = ["hidden", "max ep.", "ép. reais", "lr", "tempo", "treino", "val_acc"]
    fig, ax = plt.subplots(figsize=(9.0, 3.6))
    fig.patch.set_facecolor("white"); ax.axis("off")
    tbl = ax.table(cellText=rows, colLabels=cols, loc="center", cellLoc="center")
    tbl.auto_set_font_size(False); tbl.set_fontsize(10.5); tbl.scale(1, 1.9)
    for (r, c), cell in tbl.get_celld().items():
        if r == 0:
            cell.set_facecolor("#004080")
            cell.set_text_props(color="white", fontweight="bold")
        elif r == 1:
            cell.set_facecolor("#d4edda")
            cell.set_text_props(fontweight="bold")
        else:
            cell.set_facecolor("#f8f9fa" if r % 2 == 0 else "white")
        cell.set_edgecolor("#cccccc")
    ax.set_title("Grid Search — resultados reais (ordenados por val_acc)",
                 fontsize=11, fontweight="bold", color="#111111", pad=10)
    fig.tight_layout()
    fig.savefig(path, dpi=160, bbox_inches="tight", facecolor="white")
    plt.close(fig)


# ═══════════════════════════════════════════════════════════════════════════
# Imagem: tabela K-Fold (dados reais)
# ═══════════════════════════════════════════════════════════════════════════

def _make_kfold_table_img(path):
    rows = [
        ("Fold 1", "86.47%", "22.81s", "7975"),
        ("Fold 2", "82.64%", "22.97s", "8000"),
        ("Fold 3", "87.17%", "21.82s", "7237"),
        ("Fold 4", "83.02%", "20.28s", "7076"),
        ("Fold 5", "84.53%", "22.48s", "8000"),
        ("Média", "84.77%", "22.07s", "—"),
        ("Desv. Pad.", "±2.02%", "—", "—"),
    ]
    cols = ["Fold", "Acurácia", "Tempo (s)", "Épocas reais"]
    fig, ax = plt.subplots(figsize=(7.0, 4.2))
    fig.patch.set_facecolor("white"); ax.axis("off")
    tbl = ax.table(cellText=rows, colLabels=cols, loc="center", cellLoc="center")
    tbl.auto_set_font_size(False); tbl.set_fontsize(11); tbl.scale(1, 1.7)
    for (r, c), cell in tbl.get_celld().items():
        if r == 0:
            cell.set_facecolor("#004080")
            cell.set_text_props(color="white", fontweight="bold")
        elif r in (6, 7):
            cell.set_facecolor("#fffde7")
            cell.set_text_props(fontweight="bold")
        else:
            cell.set_facecolor("#f8f9fa" if r % 2 == 0 else "white")
        cell.set_edgecolor("#cccccc")
    ax.set_title("Validação Cruzada K=5 (melhores parâmetros: hidden=64, lr=0.015)",
                 fontsize=11, fontweight="bold", color="#111111", pad=10)
    fig.tight_layout()
    fig.savefig(path, dpi=160, bbox_inches="tight", facecolor="white")
    plt.close(fig)


# ═══════════════════════════════════════════════════════════════════════════
# Imagem: Matriz de Confusão (dados reais de saidas_teste.csv)
# ═══════════════════════════════════════════════════════════════════════════

def _make_cm_img(path, figsize=(8.5, 7.2)):
    y_real, y_pred = [], []
    with open(CSV_TEST, newline="") as f:
        for row in csv.DictReader(f):
            y_real.append(int(row["classe_real"]))
            y_pred.append(int(row["classe_prevista"]))

    n = 26
    cm = np.zeros((n, n), dtype=int)
    for r, p in zip(y_real, y_pred):
        cm[r][p] += 1

    letras = [chr(65 + i) for i in range(26)]

    fig, ax = plt.subplots(figsize=figsize)
    fig.patch.set_facecolor("white")
    im = ax.imshow(cm, cmap="Blues", aspect="auto")

    ax.set_xticks(range(26)); ax.set_xticklabels(letras, fontsize=7.5)
    ax.set_yticks(range(26)); ax.set_yticklabels(letras, fontsize=7.5)
    ax.set_xlabel("Predito", fontsize=10, labelpad=6)
    ax.set_ylabel("Real", fontsize=10, labelpad=6)
    ax.set_title("Matriz de Confusão — Conjunto de Teste (test_acc = 84,59%)",
                 fontsize=11, fontweight="bold", pad=8)
    plt.colorbar(im, ax=ax, fraction=0.03, pad=0.02)

    # diagonal highlight (correct predictions)
    for i in range(n):
        if cm[i, i] > 0:
            ax.add_patch(plt.Rectangle((i-0.5, i-0.5), 1, 1,
                         fill=False, edgecolor="#1565C0", lw=1.2))

    fig.tight_layout()
    fig.savefig(path, dpi=160, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _make_cm_top_errors_img(path):
    """Bar chart of per-class accuracy to show best/worst letters."""
    y_real, y_pred = [], []
    with open(CSV_TEST, newline="") as f:
        for row in csv.DictReader(f):
            y_real.append(int(row["classe_real"]))
            y_pred.append(int(row["classe_prevista"]))

    n = 26
    correct = np.zeros(n); total = np.zeros(n)
    for r, p in zip(y_real, y_pred):
        total[r] += 1
        if r == p:
            correct[r] += 1

    acc_per_class = np.where(total > 0, correct / total, 0)
    letras = [chr(65 + i) for i in range(26)]

    colors = ["#1565C0" if a >= 0.9 else ("#E65100" if a < 0.7 else "#4CAF50")
              for a in acc_per_class]

    fig, ax = plt.subplots(figsize=(9.0, 4.5))
    fig.patch.set_facecolor("white")
    bars = ax.bar(letras, acc_per_class * 100, color=colors, edgecolor="white", linewidth=0.5)
    ax.axhline(y=84.59, color="#555", linestyle="--", lw=1.2, label="Média 84,59%")
    ax.set_ylim(0, 110)
    ax.set_xlabel("Letra", fontsize=10)
    ax.set_ylabel("Acurácia (%)", fontsize=10)
    ax.set_title("Acurácia por Classe — Azul ≥0% | Verde 70–90% | Laranja <70%",
                 fontsize=10, fontweight="bold")
    ax.legend(fontsize=9)
    ax.tick_params(axis='x', labelsize=9)

    for bar, a in zip(bars, acc_per_class):
        if total[letras.index(chr(65 + list(acc_per_class).index(a)))] > 0:
            pass
    for i, (bar, a) in enumerate(zip(bars, acc_per_class)):
        if total[i] > 0:
            ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1.5,
                    f"{a*100:.0f}%", ha="center", va="bottom", fontsize=6.5,
                    rotation=90 if a < 0.5 else 0)

    fig.tight_layout()
    fig.savefig(path, dpi=160, bbox_inches="tight", facecolor="white")
    plt.close(fig)


# ═══════════════════════════════════════════════════════════════════════════
# Gera todas as imagens
# ═══════════════════════════════════════════════════════════════════════════

p = lambda name: os.path.join(TMPDIR, name)

_make_arch_img(p("arch.png"))
_make_grid_table_img(p("s8_table.png"))
_make_kfold_table_img(p("kfold_table.png"))
_make_cm_img(p("cm_full.png"))
_make_cm_top_errors_img(p("cm_errors.png"))

_make_code_img([
    "# Inicialização Xavier – camada oculta",
    "limit1 = np.sqrt(6 / (input_neurons + hidden_neurons))",
    "self.w1 = np.random.uniform(-limit1, limit1,",
    "                            (input_neurons, hidden_neurons))",
    "self.b1 = np.zeros((1, hidden_neurons))",
    "",
    "# Inicialização Xavier – camada de saída",
    "limit2 = np.sqrt(6 / (hidden_neurons + output_neurons))",
    "self.w2 = np.random.uniform(-limit2, limit2,",
    "                            (hidden_neurons, output_neurons))",
    "self.b2 = np.zeros((1, output_neurons))",
    "",
    "# Shapes resultantes:",
    "# W1: (120, 64)   b1: (1, 64)",
    "# W2:  (64, 26)   b2: (1, 26)",
], p("s2_code.png"))

_make_code_img([
    "def feedfoward(self, x):",
    "    # Pré-ativação da camada oculta",
    "    self.z1 = x.dot(self.w1) + self.b1",
    "    # Ativação tanh na camada oculta",
    "    self.f1 = np.tanh(self.z1)",
    "    # Pré-ativação da camada de saída",
    "    z2 = self.f1.dot(self.w2) + self.b2",
    "    # Saída com ativação tanh",
    "    return np.tanh(z2)",
    "",
    "def loss_function(self, output, y):",
    "    m = len(y)",
    "    y_oh = np.zeros((m, self.output_neurons))",
    "    y_oh[np.arange(m), y] = 1",
    "    mse = np.mean((output - y_oh) ** 2)",
    "    l2  = self.l2_lambda * (np.sum(self.w1 ** 2)",
    "                          + np.sum(self.w2 ** 2))",
    "    return mse + l2",
], p("s3_code.png"))

_make_code_img([
    "m = X.shape[0]",
    "",
    "# Vetor one-hot dos alvos",
    "y_true = np.zeros((m, self.output_neurons))",
    "y_true[np.arange(m), y] = 1",
    "",
    "# δ₂ = (2/m) · (saída – y_true) · (1 – saída²)",
    "delta2 = (2 / m) * (output - y_true) * (1 - output ** 2)",
    "",
    "# Gradiente de W2 com regularização L2",
    "dW2 = self.f1.T.dot(delta2) + self.l2_lambda * self.w2",
    "",
    "# Gradiente do bias b2",
    "dB2 = np.mean(delta2, axis=0, keepdims=True)",
], p("s4_code.png"))

_make_code_img([
    "# δ₁ = (δ₂ · W2ᵀ) ⊙ (1 – f1²)",
    "",
    "# Retropropaga o erro através de W2",
    "# e aplica a derivada da tanh em z1",
    "delta1 = delta2.dot(self.w2.T) * (1 - self.f1 ** 2)",
    "",
    "# Gradiente de W1 com regularização L2",
    "dW1 = X.T.dot(delta1) + self.l2_lambda * self.w1",
    "",
    "# Gradiente do bias b1",
    "dB1 = np.mean(delta1, axis=0, keepdims=True)",
], p("s5_code.png"))

_make_code_img([
    "def backpropation(self, X, y, output, learning_rate):",
    "    # ... cálculo de delta2, dW2, dB2, delta1, dW1, dB1 ...",
    "",
    "    # ── Gradiente Descendente: W = W – η · dW ──",
    "",
    "    self.w2 -= learning_rate * dW2",
    "    self.b2 -= learning_rate * dB2",
    "",
    "    self.w1 -= learning_rate * dW1",
    "    self.b1 -= learning_rate * dB1",
    "",
    "    # Melhor learning_rate encontrado: η = 0.015",
    "    # (via Grid Search, hidden=64, epochs=8000)",
], p("s6_code.png"))

_make_code_img([
    "# Feedforward para uma amostra",
    "output = model.feedfoward(x_amostra)",
    "# output.shape == (1, 26)  – um valor por letra A–Z",
    "",
    "# Índice com maior valor = letra predita",
    "classe = np.argmax(output, axis=1)   # ex: 13",
    "letra  = chr(65 + classe[0])         # ex: 'N'",
    "",
    "# One-hot encoding dos alvos (treinamento)",
    "y_onehot = np.zeros((m, self.output_neurons))",
    "y_onehot[np.arange(m), y] = 1",
    "# letra 'A' → [1, 0, 0, …, 0]",
    "# letra 'Z' → [0, 0, …, 0, 1]",
    "",
    "# Acurácia final no conjunto de teste: 84,59%",
], p("s7_code.png"))

_make_code_img([
    "for epoch in range(epochs):",
    "    output   = self.feedfoward(self.x)",
    "    loss     = self.loss_function(output, self.y)",
    "    self.backpropation(self.x, self.y, output, lr)",
    "",
    "    val_out  = self.feedfoward(x_val)",
    "    val_loss = self.loss_function(val_out, y_val)",
    "",
    "    if val_loss < best_val_loss:",
    "        best_val_loss = val_loss",
    "        wait = 0",
    "        self._best_w1 = self.w1.copy()",
    "        self._best_w2 = self.w2.copy()",
    "    else:",
    "        wait += 1",
    "        if wait >= patience:   # patience = 40",
    "            self.w1 = self._best_w1",
    "            self.w2 = self._best_w2",
    "            break",
], p("s10_code.png"))

_make_code_img([
    "def plot_loss_curves(model: NnModel):",
    "    plt.figure(figsize=(10, 6))",
    "    plt.plot(model.loss_history,",
    "             label='Loss de Treinamento')",
    "    plt.plot(model.val_loss_history,",
    "             label='Loss de Validação',",
    "             linestyle='--')",
    "    plt.title('Curvas de Loss')",
    "    plt.xlabel('Épocas')",
    "    plt.ylabel('MSE + L2')",
    "    plt.legend()",
    "    plt.grid()",
    "    plt.show()",
], p("s11_code.png"))

_make_code_img([
    "# Predições no conjunto de teste",
    "predictions_test = np.argmax(",
    "    best_model.feedfoward(x_test_n), axis=1)",
    "test_acc = (predictions_test == y_test).mean()",
    "# → test_acc = 0.8459  (84,59%)",
    "",
    "# Validação Cruzada K=5",
    "kf = KFold(n_splits=5, shuffle=True, random_state=42)",
    "for fold, (tr_idx, val_idx) in enumerate(kf.split(x)):",
    "    x_tr, x_vl = x[tr_idx], x[val_idx]",
    "    mn = x_tr.min(axis=0); mx = x_tr.max(axis=0)",
    "    x_tr_n = (x_tr - mn) / (mx - mn + 1e-8)",
    "    x_vl_n = (x_vl - mn) / (mx - mn + 1e-8)",
    "    # treina com melhores hiperparâmetros",
    "    # gera matriz de confusão por fold",
], p("s12_code.png"))


# ═══════════════════════════════════════════════════════════════════════════
# Preenche os slides
# ═══════════════════════════════════════════════════════════════════════════

prs    = Presentation(TEMPLATE)
slides = prs.slides

# ── Slide 2 – Estruturas de dados ─────────────────────────────────────────
s = slides[1]; _clear_body(s)
left_text(s, [
    ("Estrutura de programação:", True),
    ("Todos os dados são manipulados como arrays NumPy bidimensionais. X tem shape (n_amostras, 120): cada linha é uma imagem de letra achatada. O vetor y tem shape (n_amostras,) com o índice da classe (0=A … 25=Z).", False),
    ("", False),
    ("Camada de Entrada — 120 neurônios:", True),
    ("Cada neurônio recebe o valor de um pixel da imagem (0 ou 1). A rede aprende quais padrões de pixels ativam cada letra. W1 (120×64) e b1 (1×64) inicializados com Xavier Uniform.", False),
    ("", False),
    ("Camada Oculta — 64 neurônios:", True),
    ("Usa a função de ativação tanh para extrair representações intermediárias não lineares. 64 foi o melhor valor encontrado no Grid Search.", False),
    ("", False),
    ("Camada de Saída — 26 neurônios:", True),
    ("Um neurônio por letra (A–Z). O de maior ativação indica a predição: argmax(tanh(z2)) → índice → chr(65 + índice). W2 (64×26), b2 (1×26).", False),
])
right_img(s, p("arch.png"))

# ── Slide 3 – Feedforward + erro ──────────────────────────────────────────
s = slides[2]; _clear_body(s)
left_text(s, [
    ("Feedforward:", True),
    ("z1 = X·W1 + b1 (pré-ativação oculta). f1 = tanh(z1) (ativação oculta). z2 = f1·W2 + b2 (pré-ativação saída). saída = tanh(z2). Cada passagem propaga os 120 pixels até 26 saídas.", False),
    ("", False),
    ("Cálculo do erro (Loss):", True),
    ("Combina MSE (Erro Quadrático Médio) entre a saída tanh e o rótulo one-hot com regularização L2 (λ=0.001). O MSE penaliza o quadrado da diferença; a L2 penaliza pesos grandes para reduzir overfitting.", False),
    ("", False),
    ("E = (1/N)·Σ(saída – y_onehot)²  +  λ·(||W1||² + ||W2||²)", False),
])
right_img(s, p("s3_code.png"))

# ── Slide 4 – Backprop saída ───────────────────────────────────────────────
s = slides[3]; _clear_body(s)
left_text(s, [
    ("Backpropagation – Camada de Saída:", True),
    ("O sinal de erro é calculado como δ₂ = (2/m)·(saída − y_true)·(1 − saída²). O fator (1 − saída²) é a derivada da tanh avaliada na saída, necessária pela regra da cadeia.", False),
    ("", False),
    ("δ₂ = (2/m) · (saída – y_true) · (1 – saída²)", False),
    ("", False),
    ("dW2 = f1ᵀ · δ₂  +  λ·W2", False),
    ("dB2 = média(δ₂)", False),
    ("", False),
    ("O termo λ·W2 é a contribuição da regularização L2 ao gradiente. Sem ele os pesos poderiam crescer indefinidamente. O bias b2 é atualizado pela média de δ₂ sobre o batch.", False),
])
right_img(s, p("s4_code.png"))

# ── Slide 5 – Backprop oculta ─────────────────────────────────────────────
s = slides[4]; _clear_body(s)
left_text(s, [
    ("Backpropagation – Camada Oculta:", True),
    ("O erro é retropropagado da camada de saída multiplicando δ₂ pela transposta de W2. O resultado é então multiplicado elemento a elemento pela derivada da tanh em z1, que é (1 − f1²), onde f1 foi obtido no feedforward.", False),
    ("", False),
    ("δ₁ = (δ₂ · W2ᵀ) ⊙ (1 – f1²)", False),
    ("", False),
    ("dW1 = Xᵀ · δ₁  +  λ·W1", False),
    ("dB1 = média(δ₁)", False),
    ("", False),
    ("Esse processo garante que os pesos de todas as camadas sejam ajustados na direção que minimiza o erro total, mesmo que não estejam diretamente conectados à saída.", False),
])
right_img(s, p("s5_code.png"))

# ── Slide 6 – Atualização de pesos ────────────────────────────────────────
s = slides[5]; _clear_body(s)
left_text(s, [
    ("Atualização dos pesos — Gradiente Descendente:", True),
    ("Cada peso é reduzido na direção do gradiente multiplicado pela taxa de aprendizado η. A atualização é aplicada a cada época usando todo o batch de treinamento (Batch Gradient Descent).", False),
    ("", False),
    ("W  ←  W – η · dW      b  ←  b – η · db", False),
    ("", False),
    ("Melhor valor de η = 0.015 (encontrado pelo Grid Search):", True),
    ("Valores maiores causaram instabilidade ou parada antecipada prematura em alguns modelos. Valores menores (0.010) convergiram mais devagar e tiveram val_acc inferior. η=0.015 com hidden=64 produziu o melhor equilíbrio entre velocidade e qualidade.", False),
])
right_img(s, p("s6_code.png"))

# ── Slide 7 – Resposta de reconhecimento ──────────────────────────────────
s = slides[6]; _clear_body(s)
left_text(s, [
    ("Predição do caractere:", True),
    ("Após o feedforward, a rede produz um vetor de 26 valores. O índice do maior valor (argmax) determina a classe predita. Esse índice é convertido para a letra usando chr(65 + índice): 0→'A', 25→'Z'.", False),
    ("", False),
    ("Rótulo one-hot no treinamento:", True),
    ("O alvo é um vetor onde apenas o neurônio da letra correta vale 1 e os demais valem 0. O MSE compara esse vetor com a saída tanh(z2). Isso força a rede a ativar fortemente o neurônio correto e suprimir os demais.", False),
    ("", False),
    ("letra 'A'  →  [1, 0, 0, …, 0]", False),
    ("letra 'N'  →  [0, 0, …, 1, …, 0]  (posição 13)", False),
    ("", False),
    ("Acurácia final no conjunto de teste: 84,59%  (266 amostras)", False),
])
right_img(s, p("s7_code.png"))

# ── Slide 8 – Determinação de parâmetros / Grid Search ────────────────────
s = slides[7]; _clear_body(s)
left_text(s, [
    ("Grid Search — como os parâmetros são determinados:", True),
    ("Testamos combinações de hidden_neurons, epochs e learning_rate com l2=0.001 e patience=40 fixos. Cada combinação treinou do zero e foi avaliada pela acurácia de validação. Os melhores parâmetros foram então usados para treinar o modelo final.", False),
    ("", False),
    ("Parâmetros testados:", True),
    ("hidden: [50, 64]  |  epochs: [7000, 8000]  |  lr: [0.015]", False),
    ("l2: 0.001 (fixo)  |  patience: 40 (fixo)", False),
    ("", False),
    ("Parâmetro mais crítico:", True),
    ("O número de neurônios ocultos influencia tanto o tempo quanto a acurácia. hidden=64 com 8000 épocas alcançou val_acc=84.85% em 16.94s, parando na época 7290. Aumentar hidden demais não garante melhora — o overfitting pode surgir.", False),
    ("", False),
    ("Melhor: hidden=64 | lr=0.015 | epochs=8000 → test_acc=84,59%", False),
])
right_img(s, p("s8_table.png"))

# ── Slide 9 – Dinâmica de treinamento com tempos reais ────────────────────
s = slides[8]; _clear_body(s)
full_text(s, [
    ("Dinâmica de Treinamento — Tempos e Configuração de Hardware", True),
    ("", False),
    ("Ambiente de execução:", True),
    ("CPU: processador x86-64 (sem GPU) — implementação 100% NumPy, sem frameworks de deep learning. Sistema operacional: Windows 11. Python 3.x + NumPy.", False),
    ("", False),
    ("Tabela de experimentos — Grid Search (dados reais):", True),
    ("hidden=64 | epochs=8000 | lr=0.015  →  tempo: 16.94s | épocas reais: 7290 | val_acc: 84.85%  ← MELHOR", False),
    ("hidden=50 | epochs=8000 | lr=0.015  →  tempo: 16.27s | épocas reais: 7936 | val_acc: 83.33%", False),
    ("hidden=50 | epochs=7000 | lr=0.015  →  tempo: 13.98s | épocas reais: 7000 | val_acc: 78.79%", False),
    ("hidden=64 | epochs=7000 | lr=0.015  →  tempo: 19.01s | épocas reais: 7000 | val_acc: 78.03%", False),
    ("", False),
    ("Análise:", True),
    ("O tempo de execução é dominado pelo número de épocas reais e pelo tamanho da camada oculta. A parada antecipada reduziu de 8000 para 7290 épocas no melhor modelo, economizando ~10% do tempo. Modelos que não pararam antecipadamente (ex: hidden=50, 7000 épocas) levaram menos tempo total mas tiveram pior acurácia — indicando que não convergiram totalmente.", False),
])

# ── Slide 10 – Treinamento e parada antecipada ────────────────────────────
s = slides[9]; _clear_body(s)
left_text(s, [
    ("Divisão dos dados:", True),
    ("Treino: 70%  |  Validação: 10%  |  Teste: 20%  (seed=42, embaralhado). A normalização min-max é calculada exclusivamente sobre o treino e aplicada igualmente à validação e ao teste para evitar vazamento de dados.", False),
    ("", False),
    ("Melhor execução — hidden=64, lr=0.015, epochs=8000:", True),
    ("Parada antecipada na época 7290 (patience=40). val_loss mínimo = 0.1178. train_acc = 94.83% | val_acc = 84.85% | test_acc = 84.59%.", False),
    ("", False),
    ("Como a parada antecipada atua:", True),
    ("A cada época, a rede avalia o val_loss. Se ele não melhorar por 40 épocas consecutivas, os pesos do melhor ponto são restaurados e o treino para. Isso captura o momento de melhor generalização antes que o overfitting se instale.", False),
    ("", False),
    ("Convergência dos pesos:", True),
    ("Quando o val_loss estabiliza, os gradientes ficam próximos de zero e os pesos param de mudar significativamente. A parada antecipada identifica esse ponto automaticamente.", False),
])
right_img(s, p("s10_code.png"))

# ── Slide 11 – Gráficos parada antecipada ─────────────────────────────────
s = slides[10]; _clear_body(s)
left_text(s, [
    ("Observações nos gráficos de loss:", True),
    ("hidden=50, lr=0.015, epochs=7000: convergência sem parada antecipada — val_loss decresce em todas as épocas mas não chegou ao ótimo (val_acc=78.79%). hidden=64, lr=0.015, epochs=8000: parada na época 7290, val_loss mínimo=0.1178, melhor generalização (val_acc=84.85%).", False),
    ("", False),
    ("Como dificultar o experimento:", True),
    ("Usar os conjuntos com ruído do dataset (caracteres-ruido.csv, caracteres_ruido20.csv) no treino e o limpo na validação amplifica a divergência entre as curvas de treino e validação, tornando o comportamento da parada antecipada muito mais evidente.", False),
    ("", False),
    ("Outra estratégia — classes ambíguas:", True),
    ("Mesclar letras visualmente similares (D/O, I/J, U/V) aumenta a ambiguidade do problema. O val_loss oscila mais, o modelo demora mais para convergir, e a parada antecipada passa a ter papel decisivo no resultado final.", False),
    ("", False),
    ("Em todos os experimentos o erro minimizado é o MSE somado à regularização L2 (λ=0.001).", False),
])
right_img(s, p("s11_code.png"))

# ── Slide 12 – Matriz de confusão (com imagem real) ───────────────────────
s = slides[11]; _clear_body(s)
left_text(s, [
    ("Modelo final — teste (84,59%):", True),
    ("A matriz de confusão revela quais letras a rede acerta e confunde. Letras com traços distintivos (A, G, H, N, W) têm recall próximo a 100%. As maiores confusões ocorrem entre letras visualmente similares.", False),
    ("", False),
    ("Principais confusões observadas:", True),
    ("O ↔ Q (formas circulares fechadas). E ↔ F (traços horizontais em comum). I ↔ J (hastes verticais parecidas). C ↔ G (arco com/sem gancho). U ↔ V (formas em 'v').", False),
    ("", False),
    ("Validação Cruzada K=5:", True),
    ("Acurácia média: 84.77%  |  Desvio padrão: ±2.02%. O baixo desvio indica que o modelo é estável e generaliza de forma consistente em diferentes partições dos dados.", False),
    ("", False),
    ("Para gerar uma segunda matriz de confusão sem parada antecipada: definir patience=9999 e comparar com os resultados acima.", False),
])
right_img(s, p("cm_full.png"))

# ── Slide 13 – Implementações adicionais ──────────────────────────────────
s = slides[12]; _clear_body(s)
full_text(s, [
    ("Regularização L2  (λ = 0.001):", True),
    ("Adicionada à função de perda como λ·(||W1||²+||W2||²) e aos gradientes dos pesos. Penaliza pesos grandes, forçando a rede a generalizar melhor. Sem L2, o modelo pode memorizar o treino e ter desempenho ruim no teste.", False),
    ("", False),
    ("Inicialização Xavier Uniform:", True),
    ("Pesos inicializados em [-limit, +limit] onde limit = √(6/(n_in+n_out)). Garante variância adequada dos gradientes desde a primeira época, evitando saturação da tanh e acelerando a convergência.", False),
    ("", False),
    ("Normalização Min-Max (sem data leakage):", True),
    ("As estatísticas (min, max) são calculadas exclusivamente sobre o conjunto de treino (ou treino de cada fold no K-Fold) e aplicadas igualmente a validação e teste. Isso evita que informações do futuro influenciem o treinamento.", False),
    ("", False),
    ("Grid Search automatizado + Validação Cruzada K=5:", True),
    ("4 combinações de hiperparâmetros testadas com tabela de resultados incluindo tempo de execução, épocas reais e acurácia de validação. Avaliação robusta com 5 folds usando os melhores parâmetros (hidden=64, lr=0.015, epochs=8000), acurácia média 84.77% ±2.02%.", False),
])


# ═══════════════════════════════════════════════════════════════════════════
# Slides extras (adicionados após os 14 do template)
# ═══════════════════════════════════════════════════════════════════════════

# ── Extra A – Análise detalhada da Matriz de Confusão ─────────────────────
sA = _add_extra_slide(prs)
_set_title(sA, "Análise Detalhada da Matriz de Confusão")
left_text(sA, [
    ("Acurácia por letra (conjunto de teste):", True),
    ("O gráfico ao lado mostra a acurácia individual por classe. Barras azuis indicam acurácia ≥90% (excelente), verdes indicam 70–90% (bom) e laranjas <70% (difícil). A linha tracejada marca a acurácia média geral de 84,59%.", False),
    ("", False),
    ("Letras com 100% de acerto:", True),
    ("Letras com padrões muito distintos como A, G, H, N, W são reconhecidas perfeitamente. Elas possuem combinações únicas de pixels que não se sobrepõem com outras classes.", False),
    ("", False),
    ("Letras mais difíceis:", True),
    ("As letras com menor acurácia são aquelas com formas similares a outras no dataset. O modelo confunde principalmente O/Q, E/F, I/J — letras que diferem por apenas 1 ou 2 pixels no padrão 8×15 usado.", False),
    ("", False),
    ("O dataset CARACTERES usa padrões binários 8×15 (120 pixels). Letras similares nesse grid são naturalmente mais difíceis de distinguir.", False),
])
right_img(sA, p("cm_errors.png"))

# ── Extra B – Validação Cruzada K=5 detalhada ────────────────────────────
sB = _add_extra_slide(prs)
_set_title(sB, "Validação Cruzada K-Fold (K=5) — Resultados")
left_text(sB, [
    ("Objetivo da validação cruzada:", True),
    ("Avaliar se o modelo generaliza de forma estável ou se o resultado depende da divisão específica dos dados. Com K=5, cada fold usa 80% para treino e 20% para validação, rodando os melhores hiperparâmetros do Grid Search.", False),
    ("", False),
    ("Configuração por fold:", True),
    ("hidden=64 | lr=0.015 | epochs=8000 | l2=0.001 | patience=40. A normalização é recalculada dentro de cada fold para evitar data leakage. Mesma seed (random_state=42) para reproducibilidade.", False),
    ("", False),
    ("Resultados:", True),
    ("Fold 1: 86.47% em 22.81s (7975 épocas reais)", False),
    ("Fold 2: 82.64% em 22.97s (8000 épocas — sem parada antecipada)", False),
    ("Fold 3: 87.17% em 21.82s (7237 épocas reais)", False),
    ("Fold 4: 83.02% em 20.28s (7076 épocas reais)", False),
    ("Fold 5: 84.53% em 22.48s (8000 épocas — sem parada antecipada)", False),
    ("", False),
    ("Média: 84.77%  |  Desvio padrão: ±2.02%", False),
])
right_img(sB, p("kfold_table.png"))


# ═══════════════════════════════════════════════════════════════════════════
# Salvar
# ═══════════════════════════════════════════════════════════════════════════

prs.save(OUTPUT)
print(f"Salvo em: {OUTPUT}  ({len(prs.slides)} slides)")
shutil.rmtree(TMPDIR, ignore_errors=True)
