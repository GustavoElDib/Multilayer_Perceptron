from pptx import Presentation
from pptx.dml.color import RGBColor

prs = Presentation(r"C:\Users\gagal\Downloads\trabalho_IA_pra_analise.pptx")
for i, slide in enumerate(prs.slides):
    print(f"\n{'='*60}")
    print(f"SLIDE {i+1}  ({len(slide.shapes)} shapes)")
    print(f"{'='*60}")
    for j, sh in enumerate(slide.shapes):
        pos = f"L={sh.left//9144}pt T={sh.top//9144}pt W={sh.width//9144}pt H={sh.height//9144}pt"
        if not sh.has_text_frame:
            fill_info = ""
            try:
                if sh.fill.type is not None:
                    rgb = sh.fill.fore_color.rgb
                    fill_info = f" fill=#{rgb}"
            except:
                pass
            print(f"  [{j}] type={sh.shape_type} {pos}{fill_info} (sem texto)")
            continue
        for p_idx, para in enumerate(sh.text_frame.paragraphs):
            full = para.text
            if not full.strip():
                continue
            sizes = [r.font.size for r in para.runs if r.font.size]
            bolds = [r.font.bold for r in para.runs if r.font.bold]
            colors = []
            for r in para.runs:
                try:
                    colors.append(str(r.font.color.rgb))
                except:
                    pass
            sz  = int(sizes[0]/12700) if sizes else "?"
            b   = "BOLD" if bolds else ""
            col = colors[0] if colors else "?"
            print(f"  [{j}][p{p_idx}] {pos}  sz={sz}pt {b} color={col}  >> {full[:100]}")
